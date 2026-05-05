import os
import json
from typing import Dict, Tuple, List

from pypdf import PdfReader

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


TEXT_EXTS = {".txt", ".md", ".log", ".csv", ".json"}
SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx"} | TEXT_EXTS


def read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_json_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        obj = json.load(f)
    return json.dumps(obj, ensure_ascii=False, indent=2)


def read_pdf(path: str) -> str:
    reader = PdfReader(path)
    parts = []

    for i, page in enumerate(reader.pages, start=1):
        txt = page.extract_text() or ""
        txt = txt.strip()
        if txt:
            parts.append(f"[PDF_PAGE {i}]\n{txt}")

    return "\n\n".join(parts)


def read_docx(path: str) -> str:
    if Document is None:
        raise ImportError("python-docx is not installed. Run: python -m pip install python-docx")

    doc = Document(path)
    parts: List[str] = []

    for para in doc.paragraphs:
        txt = (para.text or "").strip()
        if txt:
            parts.append(txt)

    for table_idx, table in enumerate(doc.tables, start=1):
        table_lines = []
        for row in table.rows:
            row_text = " | ".join((cell.text or "").strip() for cell in row.cells)
            if row_text.replace("|", "").strip():
                table_lines.append(row_text)

        if table_lines:
            parts.append(f"[DOCX_TABLE {table_idx}]")
            parts.extend(table_lines)

    return "\n".join(parts)


def read_pptx(path: str) -> str:
    if Presentation is None:
        raise ImportError("python-pptx is not installed. Run: python -m pip install python-pptx")

    prs = Presentation(path)
    parts: List[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    slide_parts.append(txt)

        if slide_parts:
            parts.append(f"[PPTX_SLIDE {slide_idx}]")
            parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


def read_xlsx(path: str) -> str:
    if load_workbook is None:
        raise ImportError("openpyxl is not installed. Run: python -m pip install openpyxl")

    wb = load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []

    for ws in wb.worksheets:
        parts.append(f"[XLSX_SHEET {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            line = " | ".join(vals).strip()
            if line.replace("|", "").strip():
                parts.append(line)

    return "\n".join(parts)


def load_to_text(path: str) -> Tuple[str, Dict]:
    ext = os.path.splitext(path)[1].lower()

    meta = {
        "source": os.path.basename(path),
        "path": os.path.normpath(os.path.abspath(path)),
        "ext": ext,
    }

    if ext == ".pdf":
        return read_pdf(path), meta

    if ext in {".txt", ".md", ".log", ".csv"}:
        return read_text_file(path), meta

    if ext == ".json":
        return read_json_file(path), meta

    if ext == ".docx":
        return read_docx(path), meta

    if ext == ".pptx":
        return read_pptx(path), meta

    if ext == ".xlsx":
        return read_xlsx(path), meta

    raise ValueError(f"Unsupported file type: {ext} ({path})")